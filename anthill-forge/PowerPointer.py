# This script generates a PPTX from pasted raw slide text using python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

raw_text = """
Slide 1:
BRINGING AI IN-HOUSE:
How [Company Name] Can Deploy Secure, Private AI Without Cloud Dependencies

Presented to: [Client Name]
Date: [Date]
Confidential

---

Slide 2:
Your Data + Cloud AI = Regulatory Nightmare

Data Sovereignty: HIPAA, GDPR, CCPA violations risk
Cost Spiral: $2–5 per 1,000 tokens adds up fast
Competitive Risk: Your proprietary data training competitors' models
Lock-in: Once you're in a cloud ecosystem, escape is painful

---

Slide 3:
Cloud AI vs Enterprise vs DIY

Cloud AI:
Easy to start
Data leaves your premises
Recurring costs forever

Enterprise Platforms:
On-premise
$250k+ minimum investment
6–12 month implementation

DIY:
Maximum control
Requires PhD-level expertise
No support when things break

---

Slide 4:
The 4th Path: Enterprise-Ready, On-Premise, Affordable

Anthill Spider – Collect data
Anthill Forge – Train models
Anthill Hive – Deploy securely

---

Slide 5:
Results

364 tokens/sec throughput
22-hour training
72+ hour stability
30% loss reduction
"""

def parse_slides(text):
    blocks = [b.strip() for b in text.split("---") if b.strip()]
    slides = []
    for block in blocks:
        lines = [l.strip() for l in block.splitlines() if l.strip()]
        title = lines[0]
        body = lines[1:]
        slides.append((title, body))
    return slides

prs = Presentation()

for title, body in parse_slides(raw_text):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    tf = slide.placeholders[1].text_frame
    tf.clear()

    for line in body:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.level = 1

output_path = "/mnt/data/Anthill_Template_From_Raw_Text.pptx"
prs.save(output_path)

output_path

