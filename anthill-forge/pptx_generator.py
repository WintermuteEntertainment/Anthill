# pptx_generator.py
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
import re
from typing import List, Dict, Tuple, Optional
from dataclasses import dataclass
import os

@dataclass
class SlideTheme:
    """Theme configuration for presentation"""
    primary_color: RGBColor = RGBColor(99, 102, 241)  # Indigo
    secondary_color: RGBColor = RGBColor(139, 92, 246)  # Violet
    accent_color: RGBColor = RGBColor(16, 185, 129)  # Emerald
    warning_color: RGBColor = RGBColor(245, 158, 11)  # Amber
    error_color: RGBColor = RGBColor(239, 68, 68)  # Red
    background_color: RGBColor = RGBColor(30, 30, 46)  # Dark blue-grey
    text_color: RGBColor = RGBColor(248, 250, 252)  # Almost white
    muted_color: RGBColor = RGBColor(148, 163, 184)  # Slate 400
    
    title_font_size: int = 44
    header_font_size: int = 32
    body_font_size: int = 20
    small_font_size: int = 16
    
    title_font: str = "Inter"
    body_font: str = "Inter"
    
    slide_width: Inches = Inches(13.33)  # 16:9 aspect ratio
    slide_height: Inches = Inches(7.5)

@dataclass
class SlideContent:
    """Structured content for a single slide"""
    slide_number: int
    title: str
    content_lines: List[Dict]  # List of dicts with 'type' and 'text'
    visual_note: Optional[str] = None
    has_logo: bool = False
    has_quote: bool = False
    has_table: bool = False

class PPTXGenerator:
    """Main presentation generator class"""
    
    def __init__(self, theme: Optional[SlideTheme] = None):
        self.theme = theme or SlideTheme()
        self.prs = Presentation()
        self.prs.slide_width = self.theme.slide_width
        self.prs.slide_height = self.theme.slide_height
        
    def parse_raw_text(self, raw_text: str) -> List[SlideContent]:
        """Parse raw slide text into structured SlideContent objects"""
        slides = []
        
        # Split by slide separator (line with at least 10 underscores)
        slide_sections = re.split(r'\n_{10,}\n', raw_text.strip())
        
        for i, section in enumerate(slide_sections):
            if not section.strip():
                continue
                
            lines = [line.strip() for line in section.split('\n')]
            
            # Extract slide number and title
            slide_number = i + 1
            title = ""
            
            if lines and lines[0].lower().startswith('slide'):
                # Remove "Slide X:" prefix
                title_match = re.match(r'Slide\s*\d+:?\s*(.*)', lines[0], re.IGNORECASE)
                if title_match:
                    title = title_match.group(1).strip()
                lines = lines[1:]  # Remove the slide number line
            elif lines:
                title = lines[0]
                lines = lines[1:]
            
            # Process content lines and extract visual note
            content_lines = []
            visual_note = None
            has_logo = False
            has_quote = False
            
            for line in lines:
                if line.upper().startswith('VISUAL:'):
                    visual_note = line[7:].strip()
                    continue
                    
                if '[Anthill AI Logo]' in line or '[Logo]' in line:
                    has_logo = True
                    continue
                    
                if line.startswith('"') and line.endswith('"'):
                    has_quote = True
                    
                # Parse line type
                line_type, text = self._parse_line_type(line)
                content_lines.append({'type': line_type, 'text': text})
            
            slides.append(SlideContent(
                slide_number=slide_number,
                title=title if title else f"Slide {slide_number}",
                content_lines=content_lines,
                visual_note=visual_note,
                has_logo=has_logo,
                has_quote=has_quote
            ))
        
        return slides
    
    def _parse_line_type(self, line: str) -> Tuple[str, str]:
        """Determine the type of content line"""
        # Headers
        if line.upper() == line and line.endswith(':'):
            return 'header', line
        if line.upper() == line and len(line) > 3 and not line.endswith(':'):
            return 'subheader', line
        
        # Lists
        if line.startswith('•') or line.startswith('-') or line.startswith('*'):
            return 'bullet', line[1:].strip()
        if line.startswith('✓') or line.startswith('✅'):
            return 'check_positive', line[1:].strip()
        if line.startswith('✗') or line.startswith('❌'):
            return 'check_negative', line[1:].strip()
        if re.match(r'^\d+\.\s', line):
            return 'number', line
        
        # Icons
        icon_patterns = [
            ('🔐', 'lock'), ('💸', 'money'), ('🕵️', 'spy'), ('🔒', 'lock'),
            ('🕷️', 'spider'), ('⚒️', 'hammer'), ('🏭', 'factory'),
            ('📊', 'chart'), ('⏰', 'clock'), ('📄', 'document')
        ]
        for icon, icon_type in icon_patterns:
            if icon in line:
                return f'icon_{icon_type}', line
        
        # Regular text
        return 'text', line
    
    def build_slide_model(self, slide: SlideContent) -> Dict:
        """Build a model for slide rendering with layout information"""
        model = {
            'slide_number': slide.slide_number,
            'title': slide.title,
            'content_blocks': [],
            'visual_note': slide.visual_note,
            'has_logo': slide.has_logo,
            'has_quote': slide.has_quote,
            'layout_type': self._determine_layout_type(slide)
        }
        
        for content in slide.content_lines:
            block = {
                'type': content['type'],
                'text': content['text'],
                'style': self._get_content_style(content['type'])
            }
            model['content_blocks'].append(block)
        
        return model
    
    def _determine_layout_type(self, slide: SlideContent) -> str:
        """Determine the best layout for the slide"""
        if slide.has_logo:
            return 'title_slide'
        if slide.has_quote:
            return 'quote_slide'
        if any('header' in block['type'] for block in slide.content_lines):
            return 'section_header'
        if len(slide.content_lines) > 10:
            return 'content_heavy'
        return 'standard'
    
    def _get_content_style(self, content_type: str) -> Dict:
        """Get style properties for content type"""
        styles = {
            'header': {'font_size': 28, 'bold': True, 'color': self.theme.primary_color},
            'subheader': {'font_size': 24, 'bold': True, 'color': self.theme.secondary_color},
            'bullet': {'font_size': 20, 'bullet': True},
            'check_positive': {'font_size': 20, 'color': self.theme.accent_color},
            'check_negative': {'font_size': 20, 'color': self.theme.error_color},
            'number': {'font_size': 20, 'numbered': True},
            'icon_lock': {'font_size': 20, 'color': self.theme.primary_color},
            'icon_money': {'font_size': 20, 'color': self.theme.warning_color},
            'icon_spy': {'font_size': 20, 'color': self.theme.error_color},
            'icon_spider': {'font_size': 20, 'color': self.theme.primary_color},
            'icon_hammer': {'font_size': 20, 'color': self.theme.secondary_color},
            'text': {'font_size': 20}
        }
        return styles.get(content_type, styles['text'])
    
    def render_pptx(self, slides: List[SlideContent], output_path: str) -> str:
        """Render all slides to a PPTX file"""
        
        # Create a blank layout for custom slides
        blank_slide_layout = self.prs.slide_layouts[6]
        
        for slide_content in slides:
            slide_model = self.build_slide_model(slide_content)
            
            if slide_model['layout_type'] == 'title_slide':
                self._render_title_slide(slide_model)
            elif slide_model['layout_type'] == 'quote_slide':
                self._render_quote_slide(slide_model)
            elif slide_model['layout_type'] == 'section_header':
                self._render_section_slide(slide_model)
            else:
                self._render_standard_slide(slide_model)
        
        # Save presentation
        self.prs.save(output_path)
        return output_path
    
    def _render_title_slide(self, model: Dict):
        """Render a title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.background_color
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = model['title']
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(self.theme.title_font_size)
        title_frame.paragraphs[0].font.color.rgb = self.theme.text_color
        title_frame.paragraphs[0].font.bold = True
        
        # Add subtitle/content
        if model['content_blocks']:
            subtitle = slide.placeholders[1]
            subtitle.text = '\n'.join([block['text'] for block in model['content_blocks']])
            subtitle_frame = subtitle.text_frame
            for paragraph in subtitle_frame.paragraphs:
                paragraph.font.size = Pt(self.theme.body_font_size)
                paragraph.font.color.rgb = self.theme.text_color
        
        # Add visual placeholder
        if model['visual_note']:
            self._add_visual_placeholder(slide, model['visual_note'])
    
    def _render_standard_slide(self, model: Dict):
        """Render a standard content slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = model['title']
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.color.rgb = self.theme.primary_color
        title_frame.paragraphs[0].font.bold = True
        
        # Content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()
        
        for block in model['content_blocks']:
            p = content_frame.add_paragraph()
            p.text = block['text']
            p.font.size = Pt(block['style']['font_size'])
            
            if 'color' in block['style']:
                p.font.color.rgb = block['style']['color']
            
            if block['style'].get('bullet'):
                p.level = 0
                p.font.size = Pt(18)
            elif block['style'].get('numbered'):
                p.level = 0
                p.font.size = Pt(18)
            
            if block['style'].get('bold'):
                p.font.bold = True
        
        # Add visual placeholder
        if model['visual_note']:
            self._add_visual_placeholder(slide, model['visual_note'], is_dark=False)
    
    def _render_quote_slide(self, model: Dict):
        """Render a slide with a quote"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Set quote-style background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.primary_color
        
        # Title in light color
        title_shape = slide.shapes.title
        title_shape.text = model['title']
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Quote content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()
        
        # Format as quote
        for block in model['content_blocks']:
            p = content_frame.add_paragraph()
            p.text = block['text']
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.italic = True
        
        # Add decorative element
        left = Inches(1)
        top = Inches(2)
        width = Inches(0.2)
        height = Inches(3)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = RGBColor(255, 255, 255)
    
    def _render_section_slide(self, model: Dict):
        """Render a section header slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Gradient background
        background = slide.background
        fill = background.fill
        fill.gradient()
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = model['title'].replace('HEADLINE:', '').strip()
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.color.rgb = self.theme.primary_color
        title_frame.paragraphs[0].font.bold = True
        
        # Content with icons
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()
        
        for block in model['content_blocks']:
            p = content_frame.add_paragraph()
            p.text = block['text']
            p.font.size = Pt(20)
            
            if 'header' in block['type']:
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = self.theme.secondary_color
        
        # Add decorative element
        self._add_section_decoration(slide)
    
    def _add_visual_placeholder(self, slide, visual_note: str, is_dark: bool = True):
        """Add a visual placeholder shape with note"""
        left = Inches(8)
        top = Inches(5)
        width = Inches(4.5)
        height = Inches(1.5)
        
        # Create placeholder box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        
        # Style the placeholder
        if is_dark:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(40, 40, 60)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 245)
        
        shape.line.color.rgb = self.theme.primary_color
        shape.line.width = Pt(2)
        
        # Add text
        text_frame = shape.text_frame
        text_frame.text = f"Visual: {visual_note}"
        text_frame.paragraphs[0].font.size = Pt(10)
        
        if is_dark:
            text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 220)
        else:
            text_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 100)
        
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_section_decoration(self, slide):
        """Add decorative elements for section slides"""
        # Add accent line
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(12)
        height = Inches(0.1)
        
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme.accent_color

# Main function for testing
def main():
    """Example usage"""
    # Read example content
    with open('example_slides.txt', 'r', encoding='utf-8') as f:
        raw_text = f.read()
    
    # Generate presentation
    generator = PPTXGenerator()
    slides = generator.parse_raw_text(raw_text)
    
    output_path = os.path.join(os.path.expanduser('~'), 'Downloads', 'Anthill_Presentation.pptx')
    result = generator.render_pptx(slides, output_path)
    
    print(f"Presentation generated: {result}")
    print(f"Total slides: {len(slides)}")

if __name__ == "__main__":
    main()
